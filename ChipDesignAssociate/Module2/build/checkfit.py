"""Estimate whether any card/textbox body overflows its container shape."""
from pptx import Presentation
from pptx.util import Emu
import sys

EMU_PT = 12700
p = Presentation(sys.argv[1] if len(sys.argv) > 1 else
                 "../Module2_Topic3_DigitalLogicDesignPrinciples.pptx")
BOTTOM = 6415000
issues = []
for i, s in enumerate(list(p.slides), 1):
    shapes = list(s.shapes)
    # collect rounded-rect "cards": autoshapes with no text
    cards = [sh for sh in shapes
             if sh.shape_type == 1 and sh.has_text_frame and not sh.text_frame.text.strip()
             and sh.height > 400000 and sh.top > 1050000]
    for c in cards:
        cb = c.top + c.height
        if cb > BOTTOM + 40000:
            issues.append((i, "card runs past bottom", c.top, c.height, ""))
        # find text boxes that start inside this card
        for sh in shapes:
            if sh.shape_type != 17 or not sh.has_text_frame or sh.top >= 6500000:
                continue
            if not (c.left - 5000 <= sh.left and sh.top >= c.top and sh.top < cb):
                continue
            tf = sh.text_frame
            txt = tf.text
            if not txt.strip():
                continue
            # estimate wrapped height
            total = 0
            for para in tf.paragraphs:
                sz = 11.0
                mono = False
                for r in para.runs:
                    if r.font.size:
                        sz = r.font.size.pt
                    if r.font.name in ("Consolas",):
                        mono = True
                t = "".join(r.text for r in para.runs)
                if not t:
                    total += sz * EMU_PT * 1.30
                    continue
                cw = sz * EMU_PT * (0.60 if mono else 0.475)
                cpl = max(1, int(sh.width / cw))
                lines = max(1, -(-len(t) // cpl))
                total += lines * sz * EMU_PT * 1.34
            avail = cb - sh.top - 60000
            if total > avail:
                issues.append((i, "text overflows card", sh.top, int(total), int(avail)))
for it in issues:
    print("slide %-3d %-24s top=%-8d need=%-8s avail=%s" % it)
print("total flagged:", len(issues))
