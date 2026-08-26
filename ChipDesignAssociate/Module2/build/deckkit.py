"""Design-system helpers for the Chip Design Associate decks.

Geometry, colours and type scale are taken verbatim from
Module2_Topic1_Expanded.pptx so that Topic 3 is visually identical.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ------------------------------------------------------------------ palette
NAVY   = RGBColor(0x0E, 0x2A, 0x47)
NAVY2  = RGBColor(0x16, 0x39, 0x5E)
TEAL   = RGBColor(0x1B, 0x9A, 0xAA)
BODY   = RGBColor(0x33, 0x41, 0x4F)
INK    = RGBColor(0x1A, 0x23, 0x32)
SLATE  = RGBColor(0x5A, 0x6B, 0x7B)
GREEN  = RGBColor(0x2A, 0x9D, 0x5C)
AMBER  = RGBColor(0xC7, 0x75, 0x14)
RED    = RGBColor(0xC0, 0x1F, 0x43)
VIOLET = RGBColor(0x6B, 0x45, 0xA8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRID   = RGBColor(0xD8, 0xDE, 0xE5)
CARD   = RGBColor(0xF4, 0xF8, 0xFB)
CARD_G = RGBColor(0xEE, 0xF7, 0xF1)
CARD_A = RGBColor(0xFF, 0xF7, 0xEC)
CARD_R = RGBColor(0xFD, 0xEC, 0xEF)
PALE   = RGBColor(0xC9, 0xD6, 0xE3)
PALE2  = RGBColor(0xE7, 0xEE, 0xF5)
CODEBG = RGBColor(0x11, 0x22, 0x33)
CODEFG = RGBColor(0xD8, 0xE4, 0xF0)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

# ------------------------------------------------------------------ geometry
SW, SH   = 12161520, 6858000          # slide 13.3" x 7.5"
ML       = 411480                     # content left margin
MW       = 11247120                   # content width
TOP      = 1143000                    # first content row
BOTTOM   = 6415000                    # last usable y (footer sits below)
FOOT_Y   = 6547104
IMG_DIR  = os.environ.get("CDA_IMG_DIR",
                          os.path.join(os.path.dirname(os.path.abspath(__file__)), "img"))


def _norm(v):
    return int(round(v))


class Deck:
    def __init__(self, module_footer, code_footer):
        self.prs = Presentation()
        self.prs.slide_width = Emu(SW)
        self.prs.slide_height = Emu(SH)
        self.blank = self.prs.slide_layouts[6]
        self.module_footer = module_footer
        self.code_footer = code_footer
        self.n = 0

    # ---------------------------------------------------------- primitives
    def _shape(self, sld, kind, x, y, w, h, fill=None, line=None, lw=12700, adj=None):
        sh = sld.shapes.add_shape(kind, Emu(_norm(x)), Emu(_norm(y)),
                                  Emu(_norm(w)), Emu(_norm(max(h, 1))))
        sh.shadow.inherit = False
        if fill is None:
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = Emu(lw)
        if adj is not None:
            try:
                sh.adjustments[0] = adj
            except Exception:
                pass
        sh.text_frame.word_wrap = True
        return sh

    def rect(self, sld, x, y, w, h, fill=None, line=None, lw=12700):
        return self._shape(sld, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, lw)

    def round(self, sld, x, y, w, h, fill=None, line=None, lw=13970, adj=0.05714):
        return self._shape(sld, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                           fill, line, lw, adj)

    def ellipse(self, sld, x, y, w, h, fill=None, line=None, lw=12700):
        return self._shape(sld, MSO_SHAPE.OVAL, x, y, w, h, fill, line, lw)

    def hline(self, sld, x, y, w, color=GRID, lw=9525):
        return self.rect(sld, x, y, w, lw, fill=color, line=None)

    def text(self, sld, x, y, w, h, paras, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.03, space_after=0):
        """paras: list of paragraphs; each paragraph is a list of run dicts
        {t, b (bold), c (colour), s (size pt), f (font), i (italic)}."""
        tb = sld.shapes.add_textbox(Emu(_norm(x)), Emu(_norm(y)),
                                    Emu(_norm(w)), Emu(_norm(max(h, 1))))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        for pi, runs in enumerate(paras):
            p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            if space_after:
                p.space_after = Pt(space_after)
            if isinstance(runs, str):
                runs = [{"t": runs}]
            for rd in runs:
                r = p.add_run()
                r.text = rd.get("t", "")
                fnt = r.font
                fnt.name = rd.get("f", BODY_FONT)
                fnt.size = Pt(rd.get("s", 11))
                fnt.bold = rd.get("b", False)
                fnt.italic = rd.get("i", False)
                fnt.color.rgb = rd.get("c", BODY)
                if rd.get("spc"):
                    r.font._rPr.set("spc", str(rd["spc"]))
        return tb

    # ---------------------------------------------------------- slide types
    def title_slide(self, topic, title, subtitle, rows):
        self.n += 1
        sld = self.prs.slides.add_slide(self.blank)
        bg = sld.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        self.ellipse(sld, 8595360, -2103120, 7315200, 7315200, fill=NAVY2)
        self.ellipse(sld, 9966960, -822960, 4572000, 4572000, fill=NAVY, line=TEAL)
        self.text(sld, 640080, 1417320, 7315200, 365760,
                  [[{"t": topic, "b": True, "c": TEAL, "s": 15, "spc": 300}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, 566928, 1874520, 10424160, 1005840,
                  [[{"t": title, "b": True, "c": WHITE, "s": 35, "f": HEAD_FONT}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, 640080, 2926080, 9418320, 640080,
                  [[{"t": subtitle, "c": PALE, "s": 14}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        y = 3931920
        for r in rows:
            self.ellipse(sld, 713232, y + 82296, 237744, 237744, fill=TEAL)
            self.text(sld, 1097280, y, 9875520, 365760,
                      [[{"t": r, "c": PALE2, "s": 13.5}]], anchor=MSO_ANCHOR.MIDDLE)
            y += 457200
        self._footer(sld, dark=True)
        return sld

    def section_slide(self, kicker, title, blurb, items, accent=TEAL):
        self.n += 1
        sld = self.prs.slides.add_slide(self.blank)
        bg = sld.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        self.ellipse(sld, 9144000, -1828800, 6400800, 6400800, fill=NAVY2)
        self.text(sld, 640080, 1600200, 7315200, 365760,
                  [[{"t": kicker, "b": True, "c": accent, "s": 14, "spc": 300}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, 566928, 2011680, 9601200, 914400,
                  [[{"t": title, "b": True, "c": WHITE, "s": 30, "f": HEAD_FONT}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, 640080, 2971800, 8686800, 640080,
                  [[{"t": blurb, "c": PALE, "s": 13.5}]], anchor=MSO_ANCHOR.MIDDLE)
        y = 3840480
        for it in items:
            self.rect(sld, 640080, y + 100584, 91440, 182880, fill=accent)
            self.text(sld, 914400, y, 9875520, 365760,
                      [[{"t": it, "c": PALE2, "s": 13}]], anchor=MSO_ANCHOR.MIDDLE)
            y += 411480
        self._footer(sld, dark=True)
        return sld

    def slide(self, eyebrow, title, accent=TEAL):
        self.n += 1
        sld = self.prs.slides.add_slide(self.blank)
        self.ellipse(sld, 411480, 292608, 530352, 530352, fill=NAVY)
        self.rect(sld, 594360, 402336, 155448, 310896, fill=accent)
        self.text(sld, 1078992, 292608, 10607040, 237744,
                  [[{"t": eyebrow, "b": True, "c": accent, "s": 11, "spc": 200}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, 1060704, 502920, 10607040, 457200,
                  [[{"t": title, "b": True, "c": NAVY, "s": 21, "f": HEAD_FONT}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.hline(sld, 1078992, 1024128, 10579608)
        self._footer(sld)
        return sld

    def _footer(self, sld, dark=False):
        c = SLATE if not dark else RGBColor(0x7C, 0x8E, 0xA2)
        self.text(sld, 411480, FOOT_Y, 7315200, 237744,
                  [[{"t": self.module_footer, "c": c, "s": 8}]], anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, 8686800, FOOT_Y, 3108960, 237744,
                  [[{"t": "%s · %d" % (self.code_footer, self.n), "c": c, "s": 8}]],
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------------------------------------------------- content blocks
    def lead(self, sld, y, paras, h=None, size=12.5, w=MW, x=ML):
        n = sum(1 for _ in paras)
        if h is None:
            h = 274320 * n
        self.text(sld, x, y, w, h, paras, space_after=7)
        return y + h

    def card(self, sld, y, heading, paras, accent=TEAL, h=None, x=ML, w=MW,
             fill=CARD, size=11, head_size=11.5, pad=164592):
        if h is None:
            h = 274320 * (len(paras) + 1) + 274320
        self.round(sld, x, y, w, h, fill=fill, line=accent)
        self.text(sld, x + pad, y + 64008, w - 2 * pad, 274320,
                  [[{"t": heading, "b": True, "c": accent, "s": head_size, "f": HEAD_FONT}]],
                  anchor=MSO_ANCHOR.MIDDLE)
        self.text(sld, x + pad, y + 338328, w - 2 * pad, h - 402336,
                  paras, space_after=4)
        return y + h

    def bullets(self, sld, y, items, accent=TEAL, size=11.5, x=ML, w=MW, step=274320,
                bullet="▪"):
        for it in items:
            if isinstance(it, str):
                it = [{"t": it, "s": size}]
            self.text(sld, x, y, 182880, step,
                      [[{"t": bullet, "c": accent, "s": size * 0.8, "b": True}]],
                      anchor=MSO_ANCHOR.MIDDLE)
            self.text(sld, x + 228600, y, w - 228600, step, [it],
                      anchor=MSO_ANCHOR.MIDDLE)
            y += step
        return y

    def image(self, sld, y, name, h, w=None, x=None):
        p = os.path.join(os.environ.get("CDA_IMG_DIR", IMG_DIR), name + ".png")
        iw, ih = Image.open(p).size
        ar = iw / ih
        if w is None:
            w = MW
        dw, dh = w, w / ar
        if dh > h:
            dh = h
            dw = h * ar
        if x is None:
            x = ML + (MW - dw) / 2
        sld.shapes.add_picture(p, Emu(_norm(x)), Emu(_norm(y)),
                               Emu(_norm(dw)), Emu(_norm(dh)))
        return y + dh

    def tiers(self, sld, y, rows, h=822960, gap=45720):
        """rows: list of (label, text, colour)."""
        for lab, txt, col in rows:
            self.round(sld, ML, y, MW, h, fill=CARD, line=GRID, lw=9525)
            self.round(sld, ML, y, 1737360, h, fill=col, line=col)
            self.text(sld, ML, y, 1737360, h,
                      [[{"t": lab, "b": True, "c": WHITE, "s": 12.5, "f": HEAD_FONT}]],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self.text(sld, ML + 1920240, y + 45720, MW - 2103120, h - 91440,
                      [[{"t": txt, "c": BODY, "s": 10.5}]], anchor=MSO_ANCHOR.MIDDLE)
            y += h + gap
        return y

    def cols(self, sld, y, specs, h, gap=182880):
        """specs: list of (heading, paras, accent, fill)."""
        n = len(specs)
        w = (MW - gap * (n - 1)) / n
        for i, sp in enumerate(specs):
            head, paras, accent = sp[0], sp[1], sp[2]
            fill = sp[3] if len(sp) > 3 else CARD
            x = ML + i * (w + gap)
            self.card(sld, y, head, paras, accent=accent, h=h, x=x, w=w, fill=fill,
                      pad=137160)
        return y + h

    def code(self, sld, y, lines, h=None, x=ML, w=MW, size=10.5, title=None,
             accent=TEAL):
        lh = int(size * 12700 * 1.48)
        if h is None:
            h = lh * len(lines) + 228600 + (274320 if title else 0)
        self.round(sld, x, y, w, h, fill=CODEBG, line=accent, adj=0.03)
        ty = y + 114300
        if title:
            self.text(sld, x + 182880, ty, w - 365760, 228600,
                      [[{"t": title, "b": True, "c": TEAL, "s": 10.5, "f": BODY_FONT}]],
                      anchor=MSO_ANCHOR.MIDDLE)
            ty += 274320
        paras = []
        for ln in lines:
            if isinstance(ln, str):
                paras.append([{"t": ln or " ", "c": CODEFG, "s": size, "f": MONO_FONT}])
            else:
                paras.append([{"t": r.get("t", ""), "c": r.get("c", CODEFG),
                               "s": size, "f": MONO_FONT, "b": r.get("b", False)}
                              for r in ln])
        self.text(sld, x + 182880, ty, w - 365760, h - (ty - y) - 114300,
                  paras, line_spacing=1.12)
        return y + h

    def table(self, sld, y, headers, rows, widths, rh=228600, x=ML,
              head_fill=NAVY, size=10, head_size=10, bold_cols=(), col_colors=None):
        tw = sum(widths)
        cx = x
        for i, hcell in enumerate(headers):
            self.rect(sld, cx, y, widths[i], rh, fill=head_fill, line=head_fill, lw=6350)
            self.text(sld, cx + 45720, y, widths[i] - 91440, rh,
                      [[{"t": str(hcell), "b": True, "c": WHITE, "s": head_size}]],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += widths[i]
        yy = y + rh
        for ri, row in enumerate(rows):
            cx = x
            bg = WHITE if ri % 2 == 0 else CARD
            for i, cell in enumerate(row):
                self.rect(sld, cx, yy, widths[i], rh, fill=bg, line=GRID, lw=6350)
                col = BODY
                if col_colors and i in col_colors:
                    col = col_colors[i]
                self.text(sld, cx + 45720, yy, widths[i] - 91440, rh,
                          [[{"t": str(cell), "b": i in bold_cols,
                             "c": NAVY if i in bold_cols and not (col_colors and i in col_colors) else col,
                             "s": size}]],
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                cx += widths[i]
            yy += rh
        return yy

    def save(self, path):
        self.prs.save(path)
        print("saved", path, "-", len(self.prs.slides.__iter__.__self__._sldIdLst), "slides"
              if False else "slides:", self.n)
